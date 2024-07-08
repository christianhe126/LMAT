import base64
import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
import re
from bs4 import BeautifulSoup
from io import BytesIO
from PIL import Image

def parse_and_add_text(text_frame, html_text):
    # Define regular expressions for different HTML-like tags
    bold_re = re.compile(r'<strong>(.*?)</strong>')
    italic_re = re.compile(r'<em>(.*?)</em>')
    bullet_re = re.compile(r'<ul>(.*?)</ul>', re.DOTALL)
    li_re = re.compile(r'<li>(.*?)</li>')

    # Create a single paragraph in the text frame
    p = text_frame.add_paragraph()
    
    # Function to add formatted text to the paragraph
    def add_formatted_run(paragraph, text, bold=False, italic=False):
        run = paragraph.add_run()
        run.text = text
        run.font.bold = bold
        run.font.italic = italic

    # Split the input into blocks of text and tags
    blocks = re.split(r'(<.*?>)', html_text)
    
    for block in blocks:
        if bold_re.match(block):
            bold_text = bold_re.search(block).group(1)
            add_formatted_run(p, bold_text, bold=True)
        elif italic_re.match(block):
            italic_text = italic_re.search(block).group(1)
            add_formatted_run(p, italic_text, italic=True)
        elif bullet_re.match(block):
            ul_content = bullet_re.search(block).group(1)
            for li in li_re.findall(ul_content):
                add_formatted_run(p, "\n• " + li)  # Add bullet points in the same paragraph
        elif block.startswith('<p>') and block.endswith('</p>'):
            p_text = block[3:-4]
            add_formatted_run(p, p_text)
        elif block:
            # This handles text outside of tags
            add_formatted_run(p, block)

def insertAutomations(text, automations):
    for automation in automations:
        tag = automation.get('tag')
        result = automation.get('result')
        if not result:
            result = "[placeholder]"
        if tag:
            text = text.replace(tag, result)
    return text

def createTreeStructure(entityData):
    for item in entityData:
        if item['parentNodeID'] == 'null':
            item['parentNodeID'] = None
        if item['content']['contentType_contentTypeID'] == 'T001':
            try:
                item['content']['source'] = json.loads(item['content']['source'])
            except json.JSONDecodeError:
                item['content']['source'] = {"title": item['content']['source']}

    nodes = {item['entityID']: item for item in entityData}
    roots = []

    for item in entityData:
        parent_id = item['parentNodeID']
        if parent_id is None:
            roots.append(item)
        else:
            parent = nodes[parent_id]
            if 'children' not in parent:
                parent['children'] = []
            parent['children'].append(item)
    
    return roots

def add_html_to_text_frame(text_frame, html_input):
    # Parse the HTML input
    soup = BeautifulSoup(html_input, 'html.parser')

    # Clear any existing text
    text_frame.clear()

    first_paragraph = True

    # Add paragraphs and list items to the text frame
    for element in soup:
        if element.name == 'p':
            # Check if the paragraph is only &nbsp; or empty and add an empty line
            if element.get_text().strip() == '' or element.get_text().strip() == '\xa0':
                if first_paragraph:
                    p = text_frame.paragraphs[0]
                    p.text = ''
                    first_paragraph = False
                else:
                    p = text_frame.add_paragraph()
                    p.text = ''
                continue

            if first_paragraph:
                p = text_frame.paragraphs[0]
                first_paragraph = False
            else:
                p = text_frame.add_paragraph()
            
            if element.find('strong'):
                p.text = element.get_text()
                p.font.bold = True
            elif element.find('em'):
                p.text = element.get_text()
                p.font.italic = True
            else:
                p.text = element.get_text()
                
        elif element.name == 'ul':
            for li in element.find_all('li'):
                p = text_frame.add_paragraph()
                p.text = li.get_text()
                p.level = 1
        elif element.name == 'br':
            p = text_frame.add_paragraph()
            p.text = ''
        else:
            p = text_frame.add_paragraph()
            p.text = element.get_text()
            

def add_picture_from_b64(slide, b64_string, x, y, parent_width, parent_height):
    if b64_string.startswith("data:image/"):
        b64_string = b64_string.split(",")[1]
    
    image_data = base64.b64decode(b64_string)
    
    with Image.open(BytesIO(image_data)) as img:
        img_width, img_height = img.size
    
    width_ratio = parent_width / img_width
    height_ratio = parent_height / img_height
    
    if width_ratio < height_ratio:
        new_width = parent_width
        new_height = int(img_height * width_ratio)
        y += (parent_height - new_height) / 2
    else:
        new_width = int(img_width * height_ratio)
        new_height = parent_height
        x += (parent_width - new_width) / 2
    
    image_stream = BytesIO(image_data)
    slide.shapes.add_picture(image_stream, x, y, new_width, new_height)

def write_base64_to_file(b64_string, file_path):
    binary_data = base64.b64decode(b64_string)
    with open(file_path, 'wb') as file:
        file.write(binary_data)

def process_data(tree, path):
    prs = Presentation(path)

    for root in tree:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        if not slide:
            slide = prs.slides.add_slide()

        if 'title' in root['content']['source']:
            if slide.shapes.title:
                slide.shapes.title.text = root['content']['source'].get('title')
            else:
                textbox = slide.shapes.add_textbox(242936, 323850, 6525307, 627720)
                text_frame = textbox.text_frame
                p = text_frame.add_paragraph()
                p.text = root['content']['source'].get('title')
                p.font.size = Pt(24)

        def add_content(node, x, y, width, height):
            if 'children' in node and node['content']['contentType_contentTypeID'] == 'T001':
                n = len(node['children'])
                sorted_children = sorted(node['children'], key=lambda child: child['order'])
                for (index, child) in enumerate(sorted_children):
                    if 'vertical' in node['content']['source'] and node['content']['source'].get('vertical') == 1:
                        add_content(child,
                                    x + index * (width / n),
                                    y,
                                    width / n, 
                                    height
                                    )
                    else:
                        add_content(child, 
                                    x, 
                                    y + index * (height / n), 
                                    width, 
                                    height / n
                                    )
            elif not 'children' in node:
                if node['content']['contentType_contentTypeID'] == 'T002':
                    frame = slide.shapes.add_textbox(x, y, width, height).text_frame
                    frame.word_wrap = True
                    add_html_to_text_frame(frame, node['content']['source'])
                elif node['content']['contentType_contentTypeID'] == 'T004':
                    add_picture_from_b64(slide, node['content']['source'], x, y, width, height)

        add_content(root, 242937, 1093491, 8656646, 3683504)
    
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    pptx_b64 = base64.b64encode(pptx_io.read()).decode('utf-8')
    print(pptx_b64)

if __name__ == "__main__":
    # Read JSON data from stdin
    input_data = sys.stdin.read()
    try:
        data = json.loads(input_data)
        path = "./resources/template.pptx"
        template = data['template']
        if template and template != 'null' and template != 'undefined':
            path = "./resources/template2.pptx"
            write_base64_to_file(template, path)

        entityData = json.loads(data['entityData'])
        #automationData = json.loads(data['automationData'])['value']
        tree = createTreeStructure(entityData)
        process_data(tree, path)
    except json.JSONDecodeError as e:
        print(f"Error decoding JSON: {e}", file=sys.stderr)