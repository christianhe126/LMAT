from documents.document import Document
import re
import base64
import sys
import json
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup

class PptxDocument(Document):
    def __init__(self, format, data, template=None):
        super().__init__(format, data, template)
        self.document = Presentation(self.path)
        self.max_lines = 15
        self.chars_per_line = 50
        self.paragraphs = []
        self.page_size = {
            "x": 242937,
            "y": 1093491,
            "width": 8656646,
            "height": 3683504
        }
    
    def create_page(self, title = "[No Title]"):
        if self.document.slide_layouts[6]:
            page = self.document.slides.add_slide(self.document.slide_layouts[6])
        else:
            page = self.document.slides.add_slide()
        
        self._add_title(page, title)
        return page

    def _add_title(self, page, title):
        if page.shapes.title:
            page.shapes.title.text = title
        else:
            textbox = page.shapes.add_textbox(242936, 323850, 6525307, 627720)
            text_frame = textbox.text_frame
            p = text_frame.add_paragraph()
            p.text = title
            p.font.size = Pt(24)

    def add_image(self, page, x, y, width, height, source):
        b64_string = source
        if b64_string.startswith("data:image/"):
            b64_string = b64_string.split(",")[1]
        
        image_data = base64.b64decode(b64_string)
        
        with Image.open(BytesIO(image_data)) as img:
            img_width, img_height = img.size
        
        width_ratio = width / img_width
        height_ratio = height / img_height
        
        if width_ratio < height_ratio:
            new_width = width
            new_height = int(img_height * width_ratio)
            y += (height - new_height) / 2
        else:
            new_width = int(img_width * height_ratio)
            new_height = height
            x += (width - new_width) / 2
        
        image_stream = BytesIO(image_data)
        page.shapes.add_picture(image_stream, x, y, new_width, new_height)
    
    def add_text(self, page, x, y, width, height, source):
        text_frame = page.shapes.add_textbox(x,y,width,height).text_frame
        text_frame.word_wrap = True
        text_frame.clear()
        self.html_to_element(text_frame, source)  

    def add_table(self, page, x, y, width, height, source):
        n_row, n_col, b_style, b_width, data = self._get_table_data(source)
        shape = page.shapes.add_table(n_row, n_col,x,y,width,height)
        table = shape.table

        tbl = shape._element.graphic.graphicData.tbl
        style_id_transparent = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
        style_id_black_border = '{5940675A-B579-460E-94D1-54222C63F5DA}'
        if b_style == "hidden" or b_style == "none" or b_width == 0:
            tbl[0][-1].text = style_id_transparent
        else:
            tbl[0][-1].text = style_id_black_border
        
        for i in range(n_row):
            for j in range(n_col):
                cell = table.cell(i, j)
                text_frame = cell.text_frame
                # Check if contains image
                soup = BeautifulSoup(data[i][j], 'html.parser')
                img_tag = soup.find('img')
                if img_tag:
                    cell_x, cell_y, cell_width, cell_height = self._calculate_cell_position(table, i, j)
                    self.add_image(page, cell_x, cell_y, cell_width, cell_height, img_tag.get('src'))
                else:
                    self.html_to_element(text_frame, data[i][j])

    def _get_table_data(self, html_table):
        soup = BeautifulSoup(html_table, 'html.parser')
        table = soup.find('table')
        b_style = table.get('style').split('border-style: ')[1].split(';')[0]
        b_width = table.get('border')

        # Initialize the data array and count rows and columns
        data = []
        rows = table.find_all('tr')
        n_row = len(rows)
        n_col = 0

        # Extract table data
        for row in rows:
            cells = row.find_all('td')
            if len(cells) > n_col:
                n_col = len(cells)
            data.append([cell.decode_contents() for cell in cells])

        return n_row, n_col, b_style, b_width, data
    
    def _calculate_cell_position(self, table,row, col):
        table_gf = table._graphic_frame
        x = table_gf._element.xfrm.off.x
        y = table_gf._element.xfrm.off.y

        cell_left = x
        cell_top = y
        cell_height = table.rows[0].height
        cell_width = table.columns[0].width

        for r in range(row):
            cell_height = table.rows[r].height
            cell_top += cell_height

        for c in range(col):
            cell_width = table.columns[c].width
            cell_left += cell_width
        
        cell_height = table.rows[row].height
        cell_width = table.columns[col].width

        return cell_left, cell_top, cell_width, cell_height

    def _reduce_paragraphs(self, paragraphs):
        i = 0
        while i < len(paragraphs) - 1:
            current = paragraphs[i]
            next_paragraph = paragraphs[i + 1]
            
            if current.type == "p" and next_paragraph.type == "p" and current.estimate + next_paragraph.estimate < 15:
                # Concatenate the contents and update the estimate
                current.content += next_paragraph.content
                current.estimate += next_paragraph.estimate
                # Remove the next paragraph from the list
                paragraphs.pop(i + 1)
            elif next_paragraph.type == "img" and next_paragraph.estimate < 15 - current.estimate - 1:
                next_paragraph.estimate = 15 - current.estimate - 1
                i += 1
            else:
                i += 1

    def write(self, title, paragraphs):
        # Reduce number of text elements
        self._reduce_paragraphs(paragraphs)

        # Process the paragraphs
        current_line_number = 0
        line_height = self.page_size["height"] / self.max_lines
        page = self.create_page(title=title)
        for paragraph in paragraphs:
            #print(paragraph.content, paragraph.estimate)
            if paragraph.estimate >= 15 or current_line_number + paragraph.estimate >= 15:
                page = self.create_page(title=title)
                current_line_number = 0
                
            x = int(self.page_size["x"])
            y = int(self.page_size["y"] + current_line_number*line_height)
            width = int(self.page_size["width"])
            height = int(line_height * paragraph.estimate)
            
            if paragraph.type == "p":
                self.add_text(page,x,y,width,height,paragraph.content)
                
            elif paragraph.type == "img":
                self.add_image(page,x,y,width,height,paragraph.content)
                
            elif paragraph.type == "table":
                self.add_table(page,x,y,width,height,paragraph.content)

            current_line_number += paragraph.estimate