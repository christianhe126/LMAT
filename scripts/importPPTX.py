import sys
import json
from pptx import Presentation
from pptx.util import Inches
import re

def importPPTX():
    prs = Presentation("./resources/btp.pptx")
    slides_data = []
    content_data = []
    entity_counter = 1

    for slide in prs.slides:
        slide_content = {
            #"entityID": f"E{entity_counter:03}",
            "content": [],
            "children": []
        }
        entity_counter += 1
        stack = [slide_content]

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape == slide.shapes.title:
                    slide_content["content"].append({"source": shape.text.strip()})
                else:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text:
                            ident = paragraph.level
                            new_node = {"content": [{"source": paragraph.text.strip()}], 
                                        "children": []}
                            entity_counter += 1

                            if ident > len(stack) - 1:
                                stack[-1]["children"].append(new_node)
                            else:
                                stack = stack[:ident + 1]
                                stack[-1]["children"].append(new_node)

                            stack.append(new_node)

        slides_data.append(slide_content)

    return slides_data

def generate_hierarchy(slides_data):
    entities = []
    contents = []
    entity_count = 0
    content_count = 0

    for slide in slides_data:
        parent_content_id = f"C{content_count:03d}"
        parent_entity_id = f"E{entity_count:03d}"
        
        # Add parent entity (section)
        entities.append({
            "entityID": parent_entity_id,
            "presentationID_presentationID": 1,
            "content_contentID": parent_content_id,
            "parentNodeID": "null",
            "hierarchyLevel": 0,
            "drillState": "expanded",
            "order": entity_count
        })
        contents.append({
            "contentID": parent_content_id,
            "contentType_contentTypeID": "T001",
            "source": slide["title"]
        })
        entity_count += 1
        content_count += 1

        for order, shape_text in enumerate(slide["shapes"]):
            child_content_id = f"C{content_count:03d}"
            child_entity_id = f"E{entity_count:03d}"
            entities.append({
                "entityID": child_entity_id,
                "presentationID_presentationID": 1,
                "content_contentID": child_content_id,
                "parentNodeID": parent_entity_id,
                "hierarchyLevel": 1,
                "drillState": "expanded",
                "order": order
            })
            contents.append({
                "contentID": child_content_id,
                "contentType_contentTypeID": "T002",
                "source": shape_text
            })
            entity_count += 1
            content_count += 1

    for entity in entities:
        print(entity)

    for content in contents:
        print(content)

    return entities, contents


if __name__ == "__main__":
    # Read JSON data from stdin
    input_data = sys.stdin.read()
    try:
        data = json.loads(input_data)
        slides_data = importPPTX()
        print(slides_data)
        #generate_hierarchy(slides_data)
    except json.JSONDecodeError as e:
        print(f"Error decoding JSON: {e}", file=sys.stderr)